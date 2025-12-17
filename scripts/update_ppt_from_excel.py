#!/usr/bin/env python3
"""
Update PPT from Excel - Syncs PPT tables with Excel data
Preserves ALL PPT formatting (fonts, colors, borders, cell sizes)
Only updates text values in table cells

Supports:
- Brand detail tables (per market, slides 22+)
- Brand-by-market summary tables (slides 15-18)
- Grand total text boxes (slide 3)
"""

import json
import shutil
import re
from datetime import datetime
from pathlib import Path
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# Excel column mapping (1-based)
EXCEL_FIELD_TO_COL = {
    'brand': 4,            # D
    'budget_2026': 5,      # E
    'sufficient_2026': 7,  # G
    'gap_gbp': 8,          # H
    'gap_pct': 9,          # I
    'awa': 10,             # J
    'con': 11,             # K
    'pur': 12,             # L
    'tv': 18,              # R
    'digital': 19,         # S
    'others': 22,          # V
    'long_campaigns': 30,  # AD
    'short_campaigns': 32, # AF
    'long_pct': 34,        # AH
}

# Market to Excel row ranges (inclusive)
MARKET_ROW_RANGES = {
    'KSA': (5, 11),
    'GNE': (13, 18),
    'SOUTH AFRICA': (20, 27),
    'TURKEY': (29, 34),
    'PAKISTAN': (36, 40),
    'EGYPT': (42, 44),
    'MOROCCO': (46, 46),
    'FSA': (48, 50),
    'KENYA': (52, 54),
    'NIGERIA': (58, 59),
}

# PPT data column names (after brand column for detail tables)
PPT_DATA_FIELDS = [
    'budget_2026',
    'sufficient_2026',
    'gap_gbp',
    'gap_pct',
    'awa',
    'con',
    'pur',
    'tv',
    'digital',
    'others',
    'long_campaigns',
    'short_campaigns',
    'long_pct',
]

# Brand name aliases for matching
BRAND_ALIASES = {
    'PANADOL': ['PANADOL PAIN', 'PANADOL'],
    'SENSODYNE': ['SENSODYNE'],
    'PARODONTAX': ['PARODONTAX'],
    'CENTRUM': ['CENTRUM'],
}


def normalize_brand(s):
    """Normalize brand name for matching"""
    if not s:
        return ''
    s = str(s).strip().upper()
    s = s.replace('-', ' ').replace('  ', ' ')
    aliases = {
        'GRANDPA': 'GRAND PA',
        'GRAND-PA': 'GRAND PA',
        'MED LEMON': 'MEDLEMON',
    }
    for k, v in aliases.items():
        if k in s:
            s = s.replace(k, v)
    return s.strip()


def normalize_market(s):
    """Normalize market name for matching"""
    if not s:
        return ''
    s = str(s).strip().upper()
    # Handle variations
    if 'SOUTH' in s and 'AFRICA' in s:
        return 'SOUTH AFRICA'
    return s


def extract_market_from_slide(slide):
    """Extract market name from slide text"""
    markets = ['KSA', 'GNE', 'TURKEY', 'SOUTH AFRICA', 'EGYPT', 'MOROCCO',
               'FSA', 'KENYA', 'PAKISTAN', 'NIGERIA', 'ALGERIA']

    all_text = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            all_text.append(shape.text.upper())

    full_text = ' '.join(all_text)

    for market in markets:
        if market in full_text:
            return market
    return None


def extract_brand_from_slide(slide):
    """Extract brand name from slide text for brand-by-market tables"""
    brands = ['SENSODYNE', 'PARODONTAX', 'PANADOL', 'CENTRUM']
    
    all_text = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            all_text.append(shape.text.upper())
    
    full_text = ' '.join(all_text)
    
    for brand in brands:
        if brand in full_text:
            return brand
    return None


def detect_table_format(table):
    """
    Detect table format and return format info.
    
    Returns:
        ('brand_detail', brand_col, data_start_col) - for per-market brand tables
        ('brand_by_market', data_start_col) - for brand summary by market tables
        None - if not a recognized data table
    """
    if len(table.rows) < 2:
        return None
    
    header_cells = [table.rows[0].cells[i].text.upper().strip() 
                   for i in range(min(5, len(table.rows[0].cells)))]
    
    # Check for brand-by-market table (MARKET in first column, no BRAND/CATEGORY)
    if header_cells[0] == 'MARKET' and 'BRAND' not in ' '.join(header_cells):
        return ('brand_by_market', 1)  # data starts at column 1
    
    # Check for brand detail table (has BRAND column)
    has_brand = any('BRAND' in h for h in header_cells)
    
    if has_brand:
        brand_col = None
        for i, h in enumerate(header_cells):
            if 'BRAND' in h and 'LONG' not in h:
                brand_col = i
                break
        
        if brand_col is not None:
            return ('brand_detail', brand_col, brand_col + 1)
    
    return None


def format_value(value, field_type, original_text=''):
    """Format value for PPT display"""
    if value is None:
        return '-'

    try:
        val = float(value)
    except (ValueError, TypeError):
        return str(value) if value else '-'

    if field_type == 'currency':
        if val == 0:
            return '-'
        if '£' in original_text:
            if val < 0:
                return f"-£{abs(val):,.0f}"
            return f"£{val:,.0f}"
        if val < 0:
            return f"-{abs(val):,.0f}"
        return f"{val:,.0f}"

    elif field_type == 'percentage':
        if val == 0:
            return '-'
        pct = val * 100
        if abs(pct) < 0.5:
            return '0%'
        return f"{pct:.0f}%"

    elif field_type == 'integer':
        int_val = int(val)
        if int_val == 0:
            return '-'
        return str(int_val)

    return str(value)


def format_millions(value):
    """Format value as millions for grand total boxes"""
    if value is None:
        return '0'
    try:
        val = float(value)
        return f"{val / 1_000_000:.1f}"
    except (ValueError, TypeError):
        return '0'


def get_field_type(field_name):
    """Determine field type for formatting"""
    if field_name in ['budget_2026', 'sufficient_2026', 'gap_gbp']:
        return 'currency'
    elif field_name in ['long_campaigns', 'short_campaigns']:
        return 'integer'
    else:
        return 'percentage'


def update_cell_text(cell, new_text):
    """Update cell text while preserving formatting"""
    if not cell.text_frame.paragraphs:
        return False

    para = cell.text_frame.paragraphs[0]

    if para.runs:
        para.runs[0].text = new_text
        for run in para.runs[1:]:
            run.text = ''
    else:
        para.text = new_text

    return True


def update_shape_text(shape, old_pattern, new_value):
    """Update text in a shape matching a pattern"""
    if not hasattr(shape, 'text_frame'):
        return False
    
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old_pattern in run.text:
                # Replace the number portion
                run.text = re.sub(r'[\d.]+\s*M', f'{new_value} M', run.text)
                return True
    return False


def add_updated_label(slide, timestamp):
    """Add an 'Updated' label to the top-right corner of a slide"""
    # Check if label already exists (avoid duplicates)
    for shape in slide.shapes:
        if hasattr(shape, 'text') and 'SYNCED' in shape.text:
            # Update existing label
            if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                shape.text_frame.paragraphs[0].runs[0].text = f"SYNCED {timestamp}"
            return
    
    # Create new label - top right corner
    # Standard slide is 10" x 7.5", position at top right
    left = Inches(8.0)  # From left edge
    top = Inches(0.15)   # From top
    width = Inches(1.8)
    height = Inches(0.35)
    
    # Add rounded rectangle shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    
    # Style the shape - teal background
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 133, 124)  # Haleon teal
    shape.line.fill.background()  # No border
    
    # Add text
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    
    run = p.add_run()
    run.text = f"SYNCED {timestamp}"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)  # White text
    
    # Vertical centering
    text_frame.word_wrap = False
    from pptx.enum.text import MSO_ANCHOR
    text_frame.auto_size = None
    text_frame.paragraphs[0].space_before = Pt(2)


def read_excel_data(excel_path):
    """Read all brand data from Excel"""
    wb = load_workbook(excel_path, data_only=True)
    ws = wb['2026 Sufficiency']

    data = {}  # market -> brand -> field -> value

    for market, (start_row, end_row) in MARKET_ROW_RANGES.items():
        data[market] = {}

        for row in range(start_row, end_row + 1):
            brand_cell = ws.cell(row=row, column=EXCEL_FIELD_TO_COL['brand']).value
            if not brand_cell:
                continue

            brand = normalize_brand(brand_cell)
            data[market][brand] = {'excel_row': row, 'original_brand': brand_cell}

            for field, col in EXCEL_FIELD_TO_COL.items():
                if field == 'brand':
                    continue
                value = ws.cell(row=row, column=col).value
                data[market][brand][field] = value

    wb.close()
    return data


def aggregate_by_brand(excel_data, target_brand):
    """Aggregate data for a specific brand across all markets"""
    result = {}  # market -> field -> value
    
    target_normalized = normalize_brand(target_brand)
    
    # Get brand aliases
    brand_matches = [target_normalized]
    for key, aliases in BRAND_ALIASES.items():
        if target_normalized in [normalize_brand(a) for a in aliases]:
            brand_matches = [normalize_brand(a) for a in aliases]
            break
    
    for market, brands in excel_data.items():
        for brand_name, brand_data in brands.items():
            if any(m in brand_name or brand_name in m for m in brand_matches):
                result[market] = {}
                for field in PPT_DATA_FIELDS:
                    result[market][field] = brand_data.get(field)
                break
    
    return result


def calculate_grand_totals(excel_data):
    """Calculate grand totals across all markets and brands"""
    total_budget = 0
    total_sufficient = 0
    
    for market, brands in excel_data.items():
        for brand_name, brand_data in brands.items():
            budget = brand_data.get('budget_2026') or 0
            sufficient = brand_data.get('sufficient_2026') or 0
            total_budget += budget
            total_sufficient += sufficient
    
    return {
        'budget': total_budget,
        'sufficient': total_sufficient,
    }


def update_ppt_from_excel(ppt_path, excel_path, output_dir=None):
    """Main function to update PPT from Excel data"""

    ppt_path = Path(ppt_path)
    excel_path = Path(excel_path)

    if output_dir:
        output_dir = Path(output_dir)
    else:
        output_dir = ppt_path.parent

    output_dir.mkdir(parents=True, exist_ok=True)

    # Create backup
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    backup_path = output_dir / f"BACKUP_{ppt_path.stem}_{timestamp}.pptx"
    shutil.copy2(ppt_path, backup_path)

    # Read Excel data
    excel_data = read_excel_data(excel_path)
    
    # Calculate grand totals
    grand_totals = calculate_grand_totals(excel_data)

    # Open PPT
    prs = Presentation(ppt_path)

    changes_log = []
    warnings = []
    slides_processed = 0
    cells_updated = 0
    updated_slides = {}  # slide_idx -> slide (track which slides were updated)
    
    # Label timestamp (shorter format for display)
    label_timestamp = datetime.now().strftime('%d/%m %H:%M')

    # Process each slide
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        
        # === SLIDE 3: Grand Total Boxes ===
        if slide_num == 3:
            slide_updated = False
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text = shape.text.upper()
                    if 'BRIEFED' in text and 'GBP' in text:
                        new_val = format_millions(grand_totals['budget'])
                        if update_shape_text(shape, 'GBP', new_val):
                            cells_updated += 1
                            slide_updated = True
                            changes_log.append({
                                'slide': slide_num,
                                'type': 'grand_total',
                                'field': 'budget',
                                'new_value': f"GBP {new_val} M"
                            })
                    elif 'SUFFICIENT' in text and 'GBP' in text:
                        new_val = format_millions(grand_totals['sufficient'])
                        if update_shape_text(shape, 'GBP', new_val):
                            cells_updated += 1
                            slide_updated = True
                            changes_log.append({
                                'slide': slide_num,
                                'type': 'grand_total',
                                'field': 'sufficient',
                                'new_value': f"GBP {new_val} M"
                            })
            if slide_updated:
                updated_slides[slide_idx] = slide
            continue
        
        # === SLIDES 15-18: Brand-by-Market Tables ===
        slide_brand = extract_brand_from_slide(slide)
        
        if slide_brand and slide_num < 22:
            # Get aggregated data for this brand
            brand_market_data = aggregate_by_brand(excel_data, slide_brand)
            slide_updated = False
            
            for shape in slide.shapes:
                if not shape.has_table:
                    continue
                
                table = shape.table
                table_format = detect_table_format(table)
                
                if not table_format or table_format[0] != 'brand_by_market':
                    continue
                
                data_start_col = table_format[1]
                slides_processed += 1
                
                # Process each row (skip header)
                for row_idx in range(1, len(table.rows)):
                    row = table.rows[row_idx]
                    market_text = row.cells[0].text.strip()
                    
                    if not market_text or 'TOTAL' in market_text.upper():
                        continue
                    
                    market_normalized = normalize_market(market_text)
                    
                    if market_normalized not in brand_market_data:
                        continue
                    
                    market_data = brand_market_data[market_normalized]
                    
                    # Update each data column
                    for field_idx, field_name in enumerate(PPT_DATA_FIELDS):
                        col_idx = data_start_col + field_idx
                        
                        if col_idx >= len(row.cells):
                            break
                        
                        if field_name not in market_data:
                            continue
                        
                        excel_value = market_data[field_name]
                        cell = row.cells[col_idx]
                        original_text = cell.text.strip()
                        
                        field_type = get_field_type(field_name)
                        new_text = format_value(excel_value, field_type, original_text)
                        
                        if original_text != new_text:
                            update_cell_text(cell, new_text)
                            cells_updated += 1
                            slide_updated = True
                            
                            changes_log.append({
                                'slide': slide_num,
                                'type': 'brand_by_market',
                                'brand': slide_brand,
                                'market': market_text,
                                'field': field_name,
                                'old_value': original_text,
                                'new_value': new_text,
                            })
            
            if slide_updated:
                updated_slides[slide_idx] = slide
            continue
        
        # === SLIDES 22+: Brand Detail Tables (per market) ===
        market = extract_market_from_slide(slide)

        if not market or market not in excel_data:
            continue

        market_data = excel_data[market]
        slide_updated = False

        for shape in slide.shapes:
            if not shape.has_table:
                continue

            table = shape.table
            table_format = detect_table_format(table)
            
            if not table_format or table_format[0] != 'brand_detail':
                continue
            
            _, brand_col, data_start_col = table_format
            slides_processed += 1

            for row_idx in range(1, len(table.rows)):
                row = table.rows[row_idx]

                if brand_col >= len(row.cells):
                    continue
                    
                brand_text = row.cells[brand_col].text.strip()

                if not brand_text or 'TOTAL' in brand_text.upper():
                    continue

                normalized_brand = normalize_brand(brand_text)

                excel_brand_data = None
                for excel_brand, data in market_data.items():
                    if normalize_brand(excel_brand) == normalized_brand:
                        excel_brand_data = data
                        break
                    if normalized_brand in excel_brand or excel_brand in normalized_brand:
                        excel_brand_data = data
                        break

                if not excel_brand_data:
                    warnings.append(f"Slide {slide_num}: Brand '{brand_text}' ({market}) not found in Excel")
                    continue

                for field_idx, field_name in enumerate(PPT_DATA_FIELDS):
                    col_idx = data_start_col + field_idx
                    
                    if col_idx >= len(row.cells):
                        break
                    
                    if field_name not in excel_brand_data:
                        continue

                    excel_value = excel_brand_data[field_name]
                    cell = row.cells[col_idx]
                    original_text = cell.text.strip()

                    field_type = get_field_type(field_name)
                    new_text = format_value(excel_value, field_type, original_text)

                    if original_text != new_text:
                        update_cell_text(cell, new_text)
                        cells_updated += 1
                        slide_updated = True

                        changes_log.append({
                            'slide': slide_num,
                            'type': 'brand_detail',
                            'market': market,
                            'brand': brand_text,
                            'field': field_name,
                            'old_value': original_text,
                            'new_value': new_text,
                            'excel_row': excel_brand_data.get('excel_row')
                        })
        
        if slide_updated:
            updated_slides[slide_idx] = slide

    # Add "SYNCED" labels to all updated slides
    for slide_idx, slide in updated_slides.items():
        add_updated_label(slide, label_timestamp)

    # Save updated PPT
    output_filename = f"{ppt_path.stem}_UPDATED_{timestamp}.pptx"
    output_path = output_dir / output_filename
    prs.save(output_path)

    # Generate report
    report = {
        'timestamp': datetime.now().isoformat(),
        'input_ppt': str(ppt_path),
        'input_excel': str(excel_path),
        'output_ppt': str(output_path),
        'backup_ppt': str(backup_path),
        'summary': {
            'slides_processed': slides_processed,
            'slides_labeled': len(updated_slides),
            'cells_updated': cells_updated,
            'warnings_count': len(warnings),
            'grand_totals': grand_totals,
        },
        'changes': changes_log,
        'warnings': warnings,
    }

    report_path = output_dir / f"sync_report_{timestamp}.json"
    with open(report_path, 'w') as f:
        json.dump(report, f, indent=2, default=str)

    return {
        'success': True,
        'output_ppt': output_path,
        'backup_ppt': backup_path,
        'report_path': report_path,
        'cells_updated': cells_updated,
        'warnings': warnings,
        'changes': changes_log,
    }


def main():
    """CLI entry point"""
    import argparse

    parser = argparse.ArgumentParser(description='Update PPT from Excel data')
    parser.add_argument('--ppt', required=True, help='Path to input PPT file')
    parser.add_argument('--excel', required=True, help='Path to Excel file')
    parser.add_argument('--output-dir', help='Output directory (default: same as PPT)')

    args = parser.parse_args()

    result = update_ppt_from_excel(args.ppt, args.excel, args.output_dir)

    print("\n" + "="*60)
    print("PPT UPDATE COMPLETE")
    print("="*60)
    print(f"Cells updated: {result['cells_updated']}")
    print(f"Warnings: {len(result['warnings'])}")
    print(f"\nOutput PPT: {result['output_ppt']}")
    print(f"Backup: {result['backup_ppt']}")
    print(f"Report: {result['report_path']}")

    if result['warnings']:
        print("\n--- Warnings ---")
        for w in result['warnings'][:10]:
            print(f"  - {w}")
        if len(result['warnings']) > 10:
            print(f"  ... and {len(result['warnings']) - 10} more")

    if result['changes']:
        print("\n--- Sample Changes ---")
        for c in result['changes'][:15]:
            if c.get('type') == 'grand_total':
                print(f"  Slide {c['slide']}: Grand Total {c['field']} = {c['new_value']}")
            elif c.get('type') == 'brand_by_market':
                print(f"  Slide {c['slide']}: {c['brand']} / {c['market']} / {c['field']}")
                print(f"    {c['old_value']} -> {c['new_value']}")
            else:
                print(f"  Slide {c['slide']}: {c.get('brand', 'N/A')} / {c['field']}")
                print(f"    {c['old_value']} -> {c['new_value']}")
        if len(result['changes']) > 15:
            print(f"  ... and {len(result['changes']) - 15} more")


if __name__ == '__main__':
    main()
